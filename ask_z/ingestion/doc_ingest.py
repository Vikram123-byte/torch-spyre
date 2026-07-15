"""
ask_z/ingestion/doc_ingest.py
──────────────────────────────
Ingestion pipeline for external documents: PDF, PPTX, DOCX, and plain text.

Use this to ingest onboarding guides, access-request procedures, pod setup
instructions, architecture decks, and any other material that lives *outside*
the git repository.

Supported formats
──────────────────
  .pdf   — extracted with PyMuPDF (fitz); page-aware chunking
  .pptx  — per-slide text extraction via python-pptx
  .docx  — paragraph extraction via python-docx
  .txt / .md / .rst — plain text via SentenceSplitter

Metadata injected per chunk
────────────────────────────
  file_path       — path as supplied (relative or absolute)
  source_url      — same as file_path (override with --source-url)
  doc_type        — always "external_doc"
  component_tag   — inferred from filename or --tag argument
  last_updated    — file modification date (ISO-8601)
  content_hash    — MD5 of chunk text (delta key)
  chunk_index     — position within the source file
  version         — "external" (no git SHA)
  staleness_ttl_flag — False at ingest time

Usage
──────
  # Ingest a whole folder:
  python -m ask_z.scripts.ingest_docs --dir /path/to/docs

  # Ingest one file with a custom tag and source URL:
  python -m ask_z.scripts.ingest_docs \\
      --file /path/to/pod-setup.pdf \\
      --tag pod_setup \\
      --source-url https://ibm.box.com/s/abc123
"""

from __future__ import annotations

import hashlib
import logging
import os
import regex as re
from datetime import datetime
from pathlib import Path
from typing import Iterator

log = logging.getLogger("ask_z.ingestion.doc_ingest")

# Chunk size (characters) for page/slide text before sentence-splitting.
_CHUNK_CHARS = 1200
_OVERLAP_CHARS = 150


# ── Helpers ───────────────────────────────────────────────────────────────────


def _md5(text: str) -> str:
    return hashlib.md5(text.encode("utf-8")).hexdigest()


def _mtime_iso(path: Path) -> str:
    """Return ISO-8601 date of the file's last modification."""
    try:
        ts = path.stat().st_mtime
        return datetime.fromtimestamp(ts).strftime("%Y-%m-%d")
    except OSError:
        return "unknown"


def _infer_tag(path: Path, override: str | None) -> str:
    """Infer a component tag from the filename stem or use the override."""
    if override:
        return override
    stem = path.stem.lower()
    for noise in (
        "guide",
        "doc",
        "document",
        "presentation",
        "slides",
        "v1",
        "v2",
        "final",
    ):
        stem = stem.replace(noise, "")
    stem = stem.strip("_- ").replace(" ", "_").replace("-", "_")
    return stem or path.stem


def _sliding_window(
    text: str, chunk_size: int = _CHUNK_CHARS, overlap: int = _OVERLAP_CHARS
) -> list[str]:
    """Split *text* into overlapping character-window chunks."""
    text = text.strip()
    if not text:
        return []

    chunks: list[str] = []
    start = 0
    length = len(text)

    while start < length:
        end = min(start + chunk_size, length)
        if end < length:
            for sep in (". ", "? ", "! ", "\n\n", "\n"):
                pos = text.rfind(sep, start + chunk_size // 2, end)
                if pos != -1:
                    end = pos + len(sep)
                    break
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= length:
            break
        start = end - overlap

    return chunks


def _build_chunk(
    text: str,
    idx: int,
    *,
    file_path: Path,
    source_url: str,
    component_tag: str,
    last_updated: str,
) -> dict:
    """Build a metadata dict matching the Elasticsearch schema."""
    return {
        "text": text,
        "metadata": {
            "file_path": str(file_path),
            "source_url": source_url or str(file_path),
            "doc_type": "external_doc",
            "component_tag": component_tag,
            "last_updated": last_updated,
            "git_blame_author": "",
            "version": "external",
            "chunk_index": idx,
            "content_hash": _md5(text),
            "staleness_ttl_flag": False,
        },
    }


# ── Format extractors ─────────────────────────────────────────────────────────


def _extract_pdf(path: Path) -> list[str]:
    """Extract text from a PDF using PyMuPDF. Returns one string per page."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        log.error("PyMuPDF not installed. Run: pip install pymupdf")
        return []

    pages: list[str] = []
    try:
        doc = fitz.open(str(path))
        for page in doc:
            text = page.get_text("text").strip()
            if text:
                pages.append(text)
        doc.close()
    except Exception as exc:
        log.warning("Failed to extract PDF %s: %s", path, exc)
    return pages


def _pptx_via_python_pptx(path: Path) -> list[str] | None:
    """
    Try to extract slides using python-pptx. Returns None if the file uses
    a non-standard namespace (e.g. purl.oclc.org) that python-pptx rejects.
    """
    try:
        from pptx import Presentation
    except ImportError:
        log.error("python-pptx not installed. Run: pip install python-pptx")
        return []

    _SKIP_PLACEHOLDER_TYPES = {13, 14, 15}

    def _slide_text(slide) -> tuple[str, str]:
        title = ""
        body_lines: list[str] = []
        if slide.shapes.title:
            title = (
                (slide.shapes.title.text or "")
                .replace("\x0b", " ")
                .replace("\xa0", " ")
                .strip()
            )
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                continue
            try:
                if (
                    shape.placeholder_format
                    and shape.placeholder_format.type in _SKIP_PLACEHOLDER_TYPES
                ):
                    continue
            except Exception:
                pass
            for para in shape.text_frame.paragraphs:
                line = para.text.replace("\x0b", " ").replace("\xa0", " ").strip()
                if (
                    not line
                    or line in ("Footer",)
                    or (line.isdigit() and len(line) <= 2)
                ):
                    continue
                body_lines.append(line)
        return title, "\n".join(body_lines).strip()

    try:
        prs = Presentation(str(path))
        raw: list[tuple[int, str, str]] = []
        for i, slide in enumerate(prs.slides, start=1):
            t, b = _slide_text(slide)
            raw.append((i, t, b))
        return _merge_slides(raw)
    except Exception as exc:
        log.debug("python-pptx failed (%s) — will try raw XML fallback.", exc)
        return None


def _pptx_via_raw_xml(path: Path) -> list[str]:
    """
    Fallback PPTX extractor using raw ZIP+XML parsing.
    Works on files with non-standard namespaces that python-pptx rejects
    (e.g. files using purl.oclc.org/ooxml instead of schemas.openxmlformats.org).
    """
    import zipfile
    import regex as re

    try:
        with zipfile.ZipFile(str(path)) as z:
            names = z.namelist()

            # Determine slide order from presentation.xml (sldIdLst), not rels.
            # The rels file lists slides in arbitrary order; presentation.xml has the correct sequence.
            pres_xml = z.read("ppt/presentation.xml").decode("utf-8", errors="replace")
            # Match r:id attributes in sldId elements to get display order.
            ordered_rids = re.findall(
                r'<[^>]*:sldId\b[^>]*\br:id=["\']([^"\']+)["\']', pres_xml
            )
            # Build rid→target map from rels.
            rels_xml = z.read("ppt/_rels/presentation.xml.rels").decode(
                "utf-8", errors="replace"
            )
            rid_to_target: dict[str, str] = {}
            for m in re.finditer(
                r'Id=["\']([^"\']+)["\'][^>]*Target=["\']([^"\']+)["\']', rels_xml
            ):
                rid_to_target[m.group(1)] = m.group(2)
            # Also try reversed attribute order.
            for m in re.finditer(
                r'Target=["\']([^"\']+)["\'][^>]*Id=["\']([^"\']+)["\']', rels_xml
            ):
                rid_to_target[m.group(2)] = m.group(1)

            # Build ordered slide list; fall back to numeric sort if order detection fails.
            slide_files_ordered: list[str] = []
            for rid in ordered_rids:
                target = rid_to_target.get(rid, "")
                m = re.search(r"slide\d+", target)
                if m:
                    slide_files_ordered.append(m.group())
            if not slide_files_ordered:
                # Fallback: all slide files sorted numerically.
                slide_files_ordered = sorted(
                    [
                        re.search(r"slide\d+", n).group()
                        for n in names
                        if re.search(r"ppt/slides/slide\d+\.xml$", n)
                    ],
                    key=lambda x: int(re.search(r"\d+", x).group()),
                )

            raw: list[tuple[int, str, str]] = []
            for display_num, slide_name in enumerate(slide_files_ordered, start=1):
                zip_path = f"ppt/slides/{slide_name}.xml"
                if zip_path not in names:
                    continue
                xml = z.read(zip_path).decode("utf-8", errors="replace")
                # Extract all text runs — namespace-agnostic: match any prefix:t element.
                texts = re.findall(r"<[a-z]+:t[^>]*>(.*?)</[a-z]+:t>", xml, re.DOTALL)
                clean: list[str] = []
                for t in texts:
                    t = re.sub(r"<[^>]+>", "", t).strip()
                    t = (
                        t.replace("&#x0B;", " ")
                        .replace("&#xA0;", " ")
                        .replace("\x0b", " ")
                        .replace("\xa0", " ")
                        .strip()
                    )
                    if t and t not in ("Footer",) and not (t.isdigit() and len(t) <= 2):
                        clean.append(t)
                if clean:
                    title = clean[0] if clean else ""
                    body = "\n".join(clean[1:]) if len(clean) > 1 else ""
                    raw.append((display_num, title, body))

        return _merge_slides(raw)
    except Exception as exc:
        log.warning("Raw XML PPTX extraction failed for %s: %s", path, exc)
        return []


def _merge_slides(raw: list[tuple[int, str, str]]) -> list[str]:
    """Merge title-only slides with the next slide and assemble text blocks."""
    merged: list[str] = []
    pending_title = ""
    for slide_num, title, body in raw:
        effective_title = pending_title or title
        if not body:
            pending_title = effective_title or pending_title
            continue
        parts: list[str] = []
        if effective_title:
            parts.append(f"[Slide {slide_num}] {effective_title}")
        parts.append(body)
        merged.append("\n".join(parts))
        pending_title = ""
    if pending_title:
        merged.append(f"[Section] {pending_title}")
    return merged


def _extract_pptx(path: Path) -> list[str]:
    """
    Extract text from a PowerPoint file.
    Tries python-pptx first; falls back to raw XML extraction for files
    with non-standard namespaces (e.g. Box-exported or older PPTX files).
    """
    result = _pptx_via_python_pptx(path)
    if result is None:
        log.info("Using raw XML fallback for %s", path.name)
        result = _pptx_via_raw_xml(path)
    return result


def _extract_docx(path: Path) -> list[str]:
    """Extract text from a Word document. Returns one string per section."""
    try:
        import docx
    except ImportError:
        log.error("python-docx not installed. Run: pip install python-docx")
        return []

    sections: list[str] = []
    try:
        doc = docx.Document(str(path))
        current: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if para.style.name.startswith("Heading") and current:
                sections.append("\n".join(current))
                current = [text]
            else:
                current.append(text)
        if current:
            sections.append("\n".join(current))
    except Exception as exc:
        log.warning("Failed to extract DOCX %s: %s", path, exc)
    return sections


def _extract_text(path: Path) -> list[str]:
    """Read a plain text / Markdown / RST file as a single string."""
    try:
        return [path.read_text(encoding="utf-8", errors="replace")]
    except OSError as exc:
        log.warning("Cannot read %s: %s", path, exc)
        return []


# ── Public API ─────────────────────────────────────────────────────────────────


def chunk_document(
    file_path: Path | str,
    *,
    component_tag: str | None = None,
    source_url: str = "",
) -> list[dict]:
    """
    Chunk a single external document into a list of dicts ready for embedding.

    Each dict has keys ``text`` and ``metadata`` matching the Elasticsearch schema.
    """
    path = Path(file_path).resolve()
    if not path.exists():
        log.warning("File not found: %s", path)
        return []

    suffix = path.suffix.lower()
    tag = _infer_tag(path, component_tag)
    last_updated = _mtime_iso(path)
    src_url = source_url or str(path)

    if suffix == ".pdf":
        blocks = _extract_pdf(path)
    elif suffix == ".pptx":
        blocks = _extract_pptx(path)
    elif suffix in (".docx", ".doc"):
        blocks = _extract_docx(path)
    elif suffix in (".txt", ".md", ".rst"):
        blocks = _extract_text(path)
    else:
        log.warning("Unsupported format %s — skipping %s", suffix, path)
        return []

    if not blocks:
        log.warning("No text extracted from %s", path)
        return []

    chunks: list[dict] = []
    idx = 0
    for block in blocks:
        for chunk_text in _sliding_window(block):
            chunks.append(
                _build_chunk(
                    chunk_text,
                    idx,
                    file_path=path,
                    source_url=src_url,
                    component_tag=tag,
                    last_updated=last_updated,
                )
            )
            idx += 1

    log.info("Chunked %s → %d chunks (tag=%s)", path.name, len(chunks), tag)
    return chunks


# Box folder URL pattern — /folder/<numeric_id> — appending a filename to
# this path produces a 404.  When this prefix is detected, use it as-is so
# every file in the directory links back to the Box folder root.
_BOX_FOLDER_URL_RE = re.compile(r"^https?://[^/]*box\.com/folder/\d+", re.IGNORECASE)


def chunk_directory(
    directory: Path | str,
    *,
    component_tag: str | None = None,
    source_url_prefix: str = "",
    extensions: set[str] | None = None,
    exclude_patterns: list[str] | None = None,
) -> Iterator[dict]:
    """Walk *directory* recursively and yield chunk dicts for every supported file.

    When ``source_url_prefix`` is a Box folder URL
    (``https://*.box.com/folder/<id>``), it is used as-is for every file
    rather than appended with the relative path — Box folder-path URLs are
    always 404; the correct per-file URL requires the Box file ID which is
    only available through the Box API (use ``ingest_box.py`` for that).
    """
    if extensions is None:
        extensions = {".pdf", ".pptx", ".docx", ".txt", ".md", ".rst"}
    if exclude_patterns is None:
        exclude_patterns = ["__pycache__", ".git", ".venv", "node_modules"]

    # Detect Box folder URL prefix — do not append filenames to it.
    box_folder_prefix = bool(
        source_url_prefix and _BOX_FOLDER_URL_RE.match(source_url_prefix)
    )
    if box_folder_prefix:
        log.warning(
            "source_url_prefix looks like a Box folder URL (%s). "
            "Filenames will NOT be appended — every chunk will link to the folder root. "
            "Use ingest_box.py to get per-file /file/<id> URLs.",
            source_url_prefix,
        )

    directory = Path(directory).resolve()

    for root, dirs, files in os.walk(directory):
        root_path = Path(root)
        dirs[:] = [
            d
            for d in dirs
            if not any(pat in str(root_path / d) for pat in exclude_patterns)
        ]
        for fname in sorted(files):
            fpath = root_path / fname
            if fpath.suffix.lower() not in extensions:
                continue
            if any(pat in str(fpath) for pat in exclude_patterns):
                continue
            rel = fpath.relative_to(directory)
            if box_folder_prefix:
                url = source_url_prefix  # folder root — always openable
            elif source_url_prefix:
                url = f"{source_url_prefix}/{rel}"
            else:
                url = str(fpath)
            yield from chunk_document(
                fpath, component_tag=component_tag, source_url=url
            )
